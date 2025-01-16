import os
import fitz  # PyMuPDF
from PIL import Image
from pptx import Presentation
from pptx.util import Inches
import pandas as pd


def render_pdf_pages_as_images(pdf_path, output_folder, dpi=200):
    """Rendert die Seiten einer PDF als Bilder und speichert sie."""
    pdf_document = fitz.open(pdf_path)
    image_paths = []

    # Sicherstellen, dass der Output-Ordner existiert
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    for page_number in range(len(pdf_document)):
        page = pdf_document[page_number]
        pix = page.get_pixmap(dpi=dpi)  # Seite rendern
        image_filename = f"page{page_number + 1}.png"
        image_path = os.path.join(output_folder, image_filename)
        pix.save(image_path)  # Bild speichern
        image_paths.append(image_path)

    pdf_document.close()
    return image_paths


def create_presentation(image_paths, output_ppt_path):
    """Erstellt eine PowerPoint-Präsentation und fügt Bilder ein."""
    prs = Presentation()

    # Layout für leere Folie
    slide_layout = prs.slide_layouts[5]  # Leeres Layout

    for idx, image_path in enumerate(image_paths):
        slide = prs.slides.add_slide(slide_layout)  # Neue Folie hinzufügen
        # Bildgröße und -position anpassen
        left = Inches(1)  # Von links 1 Inch
        top = Inches(1)  # Von oben 1 Inch
        slide.shapes.add_picture(image_path, left, top, width=Inches(8), height=Inches(6))

    prs.save(output_ppt_path)
    print(f"PowerPoint-Datei wurde gespeichert: {output_ppt_path}")


def extract_table_from_pdf(pdf_path):
    """Extrahiert Tabelleninhalte aus der PDF und gibt sie als DataFrame aus."""
    pdf_document = fitz.open(pdf_path)
    all_tables = []

    for page_number in range(len(pdf_document)):
        page = pdf_document[page_number]
        text = page.get_text("dict")  # Text als strukturierte Daten
        blocks = text.get("blocks", [])

        for block in blocks:
            if "lines" in block:
                # Textzeilen der Tabelle auslesen
                table_data = []
                for line in block["lines"]:
                    words = [span["text"] for span in line["spans"]]
                    table_data.append(words)

                # Prüfen, ob Daten tabellarisch wirken
                if table_data:
                    print(f"Gefundene Tabelle auf Seite {page_number + 1}:")
                    df = pd.DataFrame(table_data[1:], columns=table_data[0])  # Erste Zeile als Header
                    print(df)
                    all_tables.append(df)

    pdf_document.close()
    return all_tables


def process_pdf_to_ppt(pdf_path):
    """Hauptprozess: Rendert PDF-Seiten, erstellt PowerPoint und extrahiert Tabellen."""
    output_folder = os.path.join(os.path.dirname(pdf_path), "output_images")
    output_ppt_path = os.path.join(os.path.dirname(pdf_path), "output_presentation.pptx")

    # Schritt 1: Seiten als Bilder rendern
    print("Rendere Seiten der PDF als Bilder...")
    image_paths = render_pdf_pages_as_images(pdf_path, output_folder)
    if not image_paths:
        print("Keine Seitenbilder in der PDF gefunden.")
        return

    # Schritt 2: PowerPoint-Präsentation erstellen
    print("Erstelle PowerPoint-Präsentation...")
    create_presentation(image_paths, output_ppt_path)

    # Schritt 3: Tabellen extrahieren
    print("Extrahiere Tabellen aus der PDF...")
    tables = extract_table_from_pdf(pdf_path)
    if not tables:
        print("Keine Tabellen in der PDF gefunden.")
    else:
        print("Tabellen erfolgreich extrahiert.")
